"""
PPT 解析器 — 基于 python-pptx。

从 PowerPoint 文件中提取幻灯片文本、备注、表格，返回 Document 列表。
"""

import hashlib
import os
from typing import List, Optional

from langchain_core.documents import Document
from tqdm import tqdm

from kb.offline.scripts.file_processors.base_parser import BaseParser


class PPTParser(BaseParser):
    """PowerPoint (.pptx) 格式解析器。"""

    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def parse(
        self,
        file_path: str,
        extract_notes: bool = True,
        extract_images: bool = True,
        **kwargs,
    ) -> List[Document]:
        """
        解析 PPT 文件。

        Args:
            file_path: .pptx 文件路径
            extract_notes: 是否提取演讲者备注
            extract_images: 是否提取图片

        Returns:
            List[Document]: 每张幻灯片一个 Document
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "请安装 python-pptx: pip install python-pptx"
            )

        prs = Presentation(file_path)
        raw_docs: List[Document] = []
        images_info = []

        for slide_idx, slide in enumerate(
            tqdm(prs.slides, desc="解析 PPT 幻灯片")
        ):
            slide_text_parts = []
            slide_images = []

            for shape in slide.shapes:
                # 提取文本
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text_parts.append(text)

                # 提取表格
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    slide_text_parts.append("[表格]\n" + "\n".join(rows))

                # 提取图片
                if extract_images:
                    img_info = self._extract_shape_image(
                        shape, slide_idx, file_path
                    )
                    if img_info:
                        slide_images.append(img_info)

            # 合并幻灯片文本
            content = "\n".join(slide_text_parts)

            # 添加演讲者备注
            if extract_notes and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    content += f"\n\n[备注]\n{notes_text}"

            if content.strip():
                unique_id = hashlib.md5(
                    f"{file_path}:slide{slide_idx}".encode("utf-8")
                ).hexdigest()
                metadata = {
                    "unique_id": unique_id,
                    "source": file_path,
                    "format": "pptx",
                    "slide": slide_idx + 1,
                    "images_info": slide_images,
                }
                raw_docs.append(Document(page_content=content, metadata=metadata))

        print(f"[PPTParser] 解析完成: {file_path} → {len(raw_docs)} 张幻灯片")
        return raw_docs

    def _extract_shape_image(
        self, shape, slide_idx: int, file_path: str
    ) -> Optional[dict]:
        """尝试从 shape 中提取图片。"""
        try:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                img_ext = image.content_type.split("/")[-1]
                if img_ext == "jpeg":
                    img_ext = "jpg"

                # 保存图片
                doc_dir = os.path.dirname(file_path)
                img_dir = os.path.join(doc_dir, ".images")
                os.makedirs(img_dir, exist_ok=True)
                img_name = f"ppt_slide{slide_idx + 1}_{shape.shape_id}.{img_ext}"
                img_path = os.path.join(img_dir, img_name)

                with open(img_path, "wb") as f:
                    f.write(image.blob)

                return {
                    "image_path": img_path,
                    "page": slide_idx + 1,
                    "title": "",
                }
        except Exception:
            pass
        return None


# ---- 便利函数 ----
def load_pptx(file_path: str, **kwargs) -> List[Document]:
    """便利函数，可直接调用。"""
    return PPTParser().parse(file_path, **kwargs)
